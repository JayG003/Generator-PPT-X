"""python-pptx helpers.

This is the project's boundary against ``python-pptx``. Renderers describe
*what* they want on a slide; this module knows *how* to say it in pptx terms.
Keeping that knowledge in one place means a future change of backend, or of
python-pptx version, touches one file rather than every renderer.

All geometry is expressed in EMU via :class:`pptx.util.Emu`, never in raw
integers, so callers cannot accidentally mix units.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx import Presentation as _NewPresentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Length, Pt

from config import (
    BLANK_LAYOUT_INDEX,
    BODY_FONT_SIZE,
    BULLET_INDENT,
    CAPTION_FONT_SIZE,
    CONTENT_HEIGHT,
    CONTENT_WIDTH,
    FONT_FAMILY,
    LEFT_MARGIN,
    LINE_SPACING,
    MIN_AUTOFIT_FONT_SIZE,
    SLIDE_HEIGHT,
    SLIDE_TITLE_FONT_SIZE,
    SLIDE_TITLE_HEIGHT,
    SLIDE_WIDTH,
    SUBTEXT_COLOR,
    TABLE_ALT_ROW_BG,
    TABLE_HEADER_BG,
    TABLE_HEADER_FG,
    TABLE_ROW_BG,
    TEXT_COLOR,
    TITLE_COLOR,
    TOP_MARGIN,
)

_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "centre": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

_ANCHORS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "center": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}

BULLET_CHARACTER = "\u2022"
SUB_BULLET_CHARACTER = "\u2013"
MAX_BULLET_DEPTH = 4


class PptError(RuntimeError):
    """Raised when a slide cannot be constructed as requested."""


# ----------------------------------------------------------------------
# Presentation and slides
# ----------------------------------------------------------------------


def configure_presentation(presentation: PresentationType) -> PresentationType:
    """Apply the project slide size to a presentation."""
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    return presentation


def add_blank_slide(
    presentation: PresentationType,
    layout_index: int = BLANK_LAYOUT_INDEX,
) -> Slide:
    """Append a slide using a layout from the presentation's master."""
    layouts = presentation.slide_layouts
    if not 0 <= layout_index < len(layouts):
        raise PptError(
            f"layout index {layout_index} is out of range; "
            f"the template provides {len(layouts)} layout(s)"
        )
    return presentation.slides.add_slide(layouts[layout_index])


def set_slide_background(slide: Slide, image_path: str | Path) -> None:
    """Cover the slide with an image, pushed behind every other shape."""
    path = Path(image_path)
    if not path.is_file():
        raise PptError(f"background image not found: '{path}'")

    picture = slide.shapes.add_picture(
        str(path),
        Emu(0),
        Emu(0),
        width=SLIDE_WIDTH,
        height=SLIDE_HEIGHT,
    )
    send_to_back(slide, picture)


def send_to_back(slide: Slide, shape) -> None:
    """Reorder ``shape`` to the bottom of the slide's z-order."""
    tree = slide.shapes._spTree
    element = shape._element
    tree.remove(element)
    tree.insert(2, element)


# ----------------------------------------------------------------------
# Text
# ----------------------------------------------------------------------


def add_textbox(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    *,
    anchor: str = "top",
    word_wrap: bool = True,
):
    """Add an empty textbox with the project's default box behaviour."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = word_wrap
    frame.vertical_anchor = _ANCHORS.get(anchor.lower(), MSO_ANCHOR.TOP)
    frame.margin_left = Emu(0)
    frame.margin_right = Emu(0)
    frame.margin_top = Emu(0)
    frame.margin_bottom = Emu(0)
    return box


def style_run(
    run,
    *,
    size: Length = BODY_FONT_SIZE,
    color: RGBColor = TEXT_COLOR,
    bold: bool = False,
    italic: bool = False,
    font: str = FONT_FAMILY,
) -> None:
    """Apply typography to a single run."""
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def set_paragraph(
    paragraph,
    text: str,
    *,
    size: Length = BODY_FONT_SIZE,
    color: RGBColor = TEXT_COLOR,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    font: str = FONT_FAMILY,
    line_spacing: float = LINE_SPACING,
    space_after: Length | None = None,
    indent_level: int = 0,
) -> None:
    """Fill a paragraph with a single styled run."""
    paragraph.text = text
    paragraph.alignment = _ALIGNMENTS.get(align.lower(), PP_ALIGN.LEFT)
    paragraph.line_spacing = line_spacing
    paragraph.level = max(0, min(indent_level, MAX_BULLET_DEPTH))
    if space_after is not None:
        paragraph.space_after = space_after
    for run in paragraph.runs:
        style_run(
            run,
            size=size,
            color=color,
            bold=bold,
            italic=italic,
            font=font,
        )


def add_paragraph(text_frame, first: bool = False):
    """Return the frame's first paragraph, or append a new one."""
    if first:
        return text_frame.paragraphs[0]
    return text_frame.add_paragraph()


def add_title(
    slide: Slide,
    text: str,
    *,
    size: Length = SLIDE_TITLE_FONT_SIZE,
    color: RGBColor = TITLE_COLOR,
    align: str = "left",
    top: Length = TOP_MARGIN,
):
    """Add the standard slide heading and return its shape."""
    box = add_textbox(
        slide,
        LEFT_MARGIN,
        top,
        CONTENT_WIDTH,
        SLIDE_TITLE_HEIGHT,
        anchor="middle",
    )
    set_paragraph(
        box.text_frame.paragraphs[0],
        text,
        size=size,
        color=color,
        bold=True,
        align=align,
    )
    return box


def add_subtitle(
    slide: Slide,
    text: str,
    top: Length,
    *,
    size: Length = CAPTION_FONT_SIZE,
    color: RGBColor = SUBTEXT_COLOR,
    align: str = "left",
):
    """Add a secondary line beneath a title or figure."""
    box = add_textbox(
        slide,
        LEFT_MARGIN,
        top,
        CONTENT_WIDTH,
        Pt(size.pt * 2),
        anchor="top",
    )
    set_paragraph(
        box.text_frame.paragraphs[0],
        text,
        size=size,
        color=color,
        italic=True,
        align=align,
    )
    return box


def add_bullets(
    slide: Slide,
    items: Sequence[str | tuple[int, str]],
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    *,
    size: Length = BODY_FONT_SIZE,
    color: RGBColor = TEXT_COLOR,
    bullet: str = BULLET_CHARACTER,
    space_after: Length = Pt(8),
):
    """Add a bullet list.

    Items are either plain strings or ``(level, text)`` pairs. Bullet glyphs
    are literal characters rather than PowerPoint list formatting: pptx has no
    public API for numbering definitions, and a literal glyph renders
    identically everywhere without XML surgery that breaks between versions.
    """
    box = add_textbox(slide, left, top, width, height)
    frame = box.text_frame

    for index, item in enumerate(items):
        level, text = item if isinstance(item, tuple) else (0, item)
        level = max(0, min(int(level), MAX_BULLET_DEPTH))
        glyph = bullet if level == 0 else SUB_BULLET_CHARACTER
        paragraph = add_paragraph(frame, first=index == 0)
        set_paragraph(
            paragraph,
            f"{glyph}  {text}" if glyph else text,
            size=Pt(max(size.pt - level * 2, MIN_AUTOFIT_FONT_SIZE.pt)),
            color=color,
            space_after=space_after,
            indent_level=level,
        )
        paragraph.space_before = Emu(0)
    return box


def fit_font_size(
    text: str,
    box_width: Length,
    box_height: Length,
    *,
    base_size: Length = BODY_FONT_SIZE,
    minimum: Length = MIN_AUTOFIT_FONT_SIZE,
    line_spacing: float = LINE_SPACING,
) -> Length:
    """Estimate the largest point size at which ``text`` fits.

    An approximation, not a layout engine: PowerPoint reflows text itself and
    pptx cannot measure a rendered run. Average glyph advance is taken as
    0.5 em, which is close enough for Calibri to avoid overflow while staying
    cheap, and it degrades safely because it over-estimates width.
    """
    if not text.strip():
        return base_size

    width_pt = box_width.pt
    height_pt = box_height.pt
    characters = len(text)

    size = base_size.pt
    while size > minimum.pt:
        per_line = max(1, int(width_pt / (size * 0.5)))
        lines = -(-characters // per_line)
        if lines * size * line_spacing <= height_pt:
            return Pt(size)
        size -= 1
    return minimum


# ----------------------------------------------------------------------
# Pictures
# ----------------------------------------------------------------------


def picture_placement(
    image_width: int,
    image_height: int,
    box_left: Length,
    box_top: Length,
    box_width: Length,
    box_height: Length,
) -> tuple[Emu, Emu, Emu, Emu]:
    """Fit pixel dimensions into a box, preserving aspect ratio and centring.

    Returns ``(left, top, width, height)`` in EMU. Computed here rather than
    letting pptx infer size from DPI metadata, because matplotlib writes a DPI
    that would otherwise scale figures unpredictably between renders.
    """
    if image_width <= 0 or image_height <= 0:
        raise PptError(
            f"invalid image dimensions: {image_width}x{image_height}"
        )

    scale = min(
        box_width / image_width,
        box_height / image_height,
    )
    width = Emu(int(image_width * scale))
    height = Emu(int(image_height * scale))
    left = Emu(int(box_left + (box_width - width) / 2))
    top = Emu(int(box_top + (box_height - height) / 2))
    return left, top, width, height


def add_picture_fitted(
    slide: Slide,
    image_path: str | Path,
    image_size: tuple[int, int],
    box_left: Length,
    box_top: Length,
    box_width: Length,
    box_height: Length,
):
    """Insert an image centred inside a box, preserving its aspect ratio."""
    path = Path(image_path)
    if not path.is_file():
        raise PptError(f"image not found: '{path}'")

    left, top, width, height = picture_placement(
        image_size[0],
        image_size[1],
        box_left,
        box_top,
        box_width,
        box_height,
    )
    return slide.shapes.add_picture(
        str(path),
        left,
        top,
        width=width,
        height=height,
    )


# ----------------------------------------------------------------------
# Layout arithmetic
# ----------------------------------------------------------------------


def content_box(*, with_title: bool = True) -> tuple[Emu, Emu, Emu, Emu]:
    """The usable area of a slide as ``(left, top, width, height)``."""
    top = Emu(TOP_MARGIN + (SLIDE_TITLE_HEIGHT if with_title else 0))
    height = Emu(
        CONTENT_HEIGHT - (SLIDE_TITLE_HEIGHT if with_title else 0)
    )
    return Emu(LEFT_MARGIN), top, Emu(CONTENT_WIDTH), height


def split_columns(
    left: Length,
    width: Length,
    count: int,
    gutter: Length = Emu(228600),
) -> list[tuple[Emu, Emu]]:
    """Divide a horizontal span into ``count`` equal columns.

    Returns ``(left, width)`` per column. The default gutter is 0.25 inch.
    """
    if count < 1:
        raise PptError(f"column count must be positive, got {count}")

    total_gutter = gutter * (count - 1)
    column_width = int((width - total_gutter) / count)
    if column_width <= 0:
        raise PptError(
            f"{count} columns do not fit in the available width"
        )
    return [
        (Emu(int(left + index * (column_width + gutter))), Emu(column_width))
        for index in range(count)
    ]


def stack_rows(
    top: Length,
    height: Length,
    count: int,
    gutter: Length = Emu(114300),
) -> list[tuple[Emu, Emu]]:
    """Divide a vertical span into ``count`` equal rows."""
    if count < 1:
        raise PptError(f"row count must be positive, got {count}")

    total_gutter = gutter * (count - 1)
    row_height = int((height - total_gutter) / count)
    if row_height <= 0:
        raise PptError(f"{count} rows do not fit in the available height")
    return [
        (Emu(int(top + index * (row_height + gutter))), Emu(row_height))
        for index in range(count)
    ]


def to_rgb(color: Iterable[int] | RGBColor) -> RGBColor:
    """Coerce a ``(r, g, b)`` triple into an ``RGBColor``."""
    if isinstance(color, RGBColor):
        return color
    channels = tuple(int(value) for value in color)
    if len(channels) != 3:
        raise PptError(f"expected three colour channels, got {len(channels)}")
    return RGBColor(*channels)


# ----------------------------------------------------------------------
# Presentation construction and metadata
# ----------------------------------------------------------------------


def create_presentation() -> PresentationType:
    """A new, correctly sized, empty presentation."""
    return configure_presentation(_NewPresentation())


def set_core_properties(
    presentation: PresentationType,
    *,
    title: str = "",
    author: str = "",
    subject: str = "",
    comments: str = "",
) -> None:
    """Populate the document metadata PowerPoint shows in File > Info."""
    properties = presentation.core_properties
    if title:
        properties.title = title
    if author:
        properties.author = author
    if subject:
        properties.subject = subject
    if comments:
        properties.comments = comments


def set_background_color(slide: Slide, color: Iterable[int] | RGBColor) -> None:
    """Fill a slide with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = to_rgb(color)


# ----------------------------------------------------------------------
# Paragraph plumbing
# ----------------------------------------------------------------------


def clear_text_frame(text_frame) -> None:
    """Empty a text frame without leaving the first run's formatting behind.

    ``TextFrame.clear()`` keeps the first paragraph and its run properties,
    so the next caller silently inherits the previous font. Removing the run
    elements is the only way to start genuinely clean.
    """
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    for run in list(paragraph.runs):
        paragraph._p.remove(run._r)


def suppress_bullet(paragraph) -> None:
    """Disable the automatic bullet PowerPoint applies to body paragraphs.

    The project draws bullet characters itself so that PowerPoint, LibreOffice
    and Google Slides -- which resolve theme list styles differently -- all
    render identically.
    """
    properties = paragraph._p.get_or_add_pPr()
    if properties.find(qn("a:buNone")) is None:
        properties.append(properties.makeelement(qn("a:buNone"), {}))


# ----------------------------------------------------------------------
# Shapes
# ----------------------------------------------------------------------


def add_rectangle(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    *,
    fill: Iterable[int] | RGBColor | None = None,
    line: Iterable[int] | RGBColor | None = None,
    line_width: Length = Pt(1),
):
    """Add a plain rectangle, used for rules, accent bars and banners."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(int(left)), Emu(int(top)),
        Emu(int(width)), Emu(int(height)),
    )
    shape.shadow.inherit = False

    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = to_rgb(fill)

    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = to_rgb(line)
        shape.line.width = line_width

    if shape.has_text_frame:
        clear_text_frame(shape.text_frame)

    return shape


def add_rule(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    *,
    thickness: Length = Emu(19050),
    color: Iterable[int] | RGBColor = TITLE_COLOR,
):
    """A thin horizontal rule, used under slide titles."""
    return add_rectangle(
        slide, left, top, width, thickness, fill=color, line=None
    )


# ----------------------------------------------------------------------
# Tables
# ----------------------------------------------------------------------


def add_table(
    slide: Slide,
    rows: int,
    columns: int,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
):
    """Add an empty table with theme banding disabled.

    Banding is switched off because the renderer colours rows explicitly;
    leaving it on means the theme and the renderer disagree about alternate
    rows and the result depends on which application opens the file.
    """
    if rows < 1 or columns < 1:
        raise PptError(
            f"a table needs at least one row and column, got {rows}x{columns}"
        )

    table = slide.shapes.add_table(
        rows, columns, Emu(int(left)), Emu(int(top)),
        Emu(int(width)), Emu(int(height)),
    ).table
    table.first_row = False
    table.horz_banding = False
    return table


def set_column_widths(table, widths: Sequence[Length]) -> None:
    """Set the width of every column."""
    if len(widths) != len(table.columns):
        raise PptError(
            f"got {len(widths)} width(s) for {len(table.columns)} column(s)"
        )
    for column, width in zip(table.columns, widths):
        column.width = Emu(int(width))


def proportional_widths(
    total: Length, weights: Sequence[float]
) -> list[Emu]:
    """Split a width across columns in proportion to ``weights``."""
    if not weights:
        raise PptError("at least one column weight is required")
    if any(weight <= 0 for weight in weights):
        raise PptError("column weights must be positive")

    scale = sum(weights)
    return [Emu(int(total * weight / scale)) for weight in weights]


def set_cell(
    cell,
    text: str,
    *,
    size: Length = BODY_FONT_SIZE,
    color: Iterable[int] | RGBColor = TEXT_COLOR,
    fill: Iterable[int] | RGBColor | None = None,
    bold: bool = False,
    align: str = "left",
    anchor: str = "middle",
    font: str = FONT_FAMILY,
    margin: Length = Emu(54864),
) -> None:
    """Write and format one table cell."""
    cell.vertical_anchor = _ANCHORS.get(anchor, MSO_ANCHOR.MIDDLE)
    cell.margin_left = Emu(int(margin))
    cell.margin_right = Emu(int(margin))
    cell.margin_top = Emu(int(margin * 0.6))
    cell.margin_bottom = Emu(int(margin * 0.6))

    if fill is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = to_rgb(fill)

    frame = cell.text_frame
    clear_text_frame(frame)
    frame.word_wrap = True

    paragraph = add_paragraph(frame, first=True)
    set_paragraph(
        paragraph,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
        font=font,
        space_after=Pt(0),
    )
    suppress_bullet(paragraph)


def fill_table(
    table,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    *,
    align: str = "left",
    header_fill: Iterable[int] | RGBColor = TABLE_HEADER_BG,
    header_color: Iterable[int] | RGBColor = TABLE_HEADER_FG,
    row_fill: Iterable[int] | RGBColor = TABLE_ROW_BG,
    alt_row_fill: Iterable[int] | RGBColor = TABLE_ALT_ROW_BG,
    body_size: Length = BODY_FONT_SIZE,
    header_size: Length = BODY_FONT_SIZE,
) -> None:
    """Populate a table with an optional header row and zebra striping."""
    offset = 1 if headers else 0

    for index, heading in enumerate(headers):
        set_cell(
            table.cell(0, index), heading, size=header_size,
            color=header_color, fill=header_fill, bold=True,
            align="center" if align == "center" else "left",
        )

    for row_index, row in enumerate(rows):
        fill = alt_row_fill if row_index % 2 else row_fill
        for column_index, value in enumerate(row):
            set_cell(
                table.cell(row_index + offset, column_index), value,
                size=body_size, fill=fill, align=align,
            )


# ----------------------------------------------------------------------
# Notes
# ----------------------------------------------------------------------


def set_speaker_notes(slide: Slide, text: str) -> None:
    """Write the speaker notes for a slide, skipping empty text.

    Guarded because touching ``notes_slide`` creates the notes part in the
    package; doing that for every slide inflates a large deck for nothing.
    """
    if not text or not text.strip():
        return
    slide.notes_slide.notes_text_frame.text = text


__all__ = [
    "PptError",
    "BULLET_CHARACTER",
    "SUB_BULLET_CHARACTER",
    "configure_presentation",
    "add_blank_slide",
    "set_slide_background",
    "send_to_back",
    "add_textbox",
    "style_run",
    "set_paragraph",
    "add_paragraph",
    "add_title",
    "add_subtitle",
    "add_bullets",
    "fit_font_size",
    "picture_placement",
    "add_picture_fitted",
    "content_box",
    "split_columns",
    "stack_rows",
    "to_rgb",
    "create_presentation",
    "set_core_properties",
    "set_background_color",
    "clear_text_frame",
    "suppress_bullet",
    "add_rectangle",
    "add_rule",
    "add_table",
    "set_column_widths",
    "proportional_widths",
    "set_cell",
    "fill_table",
    "set_speaker_notes",
]